import os
import sys
import subprocess

# Check and install python-pptx if not present
try:
    import pptx
except ImportError:
    print("Installing python-pptx...")
    subprocess.run([sys.executable, "-m", "pip", "install", "python-pptx"], check=True)
    import pptx

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def apply_background(slide, color):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def create_pptx():
    prs = Presentation()
    
    # 16:9 Aspect Ratio
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    blank_slide_layout = prs.slide_layouts[6]
    
    # Color palette
    DARK_BLUE = RGBColor(11, 29, 58)      # Theme background (Dark premium)
    LIGHT_BG = RGBColor(245, 247, 250)     # Optional light background
    WHITE = RGBColor(255, 255, 255)
    GOLD = RGBColor(230, 126, 34)         # Accent color
    MUTED_GRAY = RGBColor(170, 183, 184)
    TEXT_DARK = RGBColor(44, 62, 80)
    
    slides_data = [
        # Slide 1
        {
            "bg": DARK_BLUE,
            "title": "BẢN ĐỒ KỶ NIỆM DU LỊCH\nTRAVEL MEMORY MAP",
            "title_color": GOLD,
            "title_size": 44,
            "subtitle": "Đồ án tốt nghiệp / Báo cáo môn học\nSinh viên thực hiện: Lâm Hòa Bình",
            "subtitle_color": WHITE,
            "subtitle_size": 18,
            "points": []
        },
        # Slide 2
        {
            "bg": DARK_BLUE,
            "title": "ĐẶT VẤN ĐỀ & GIẢI PHÁP",
            "title_color": GOLD,
            "title_size": 36,
            "subtitle": "",
            "subtitle_color": WHITE,
            "subtitle_size": 16,
            "points": [
                "• Rời rạc: Ảnh lưu trên điện thoại thiếu liên kết địa lý",
                "• Trực quan: Bản đồ số động kết nối tọa độ thực",
                "• Cảm xúc: Lưu giữ khoảnh khắc và lộ trình chi tiết"
            ]
        },
        # Slide 3
        {
            "bg": DARK_BLUE,
            "title": "MỤC TIÊU ĐỀ TÀI",
            "title_color": GOLD,
            "title_size": 36,
            "subtitle": "",
            "subtitle_color": WHITE,
            "subtitle_size": 16,
            "points": [
                "• Số hóa: Lưu trữ hành trình du lịch cá nhân & nhóm",
                "• Tương tác: Kết bạn, chia sẻ và tương tác trực tuyến",
                "• Thông minh: Tích hợp AI hỗ trợ gợi ý du lịch"
            ]
        },
        # Slide 4
        {
            "bg": DARK_BLUE,
            "title": "CÔNG NGHỆ SỬ DỤNG",
            "title_color": GOLD,
            "title_size": 36,
            "subtitle": "Kiến trúc tối ưu, hiện đại",
            "subtitle_color": WHITE,
            "subtitle_size": 16,
            "points": [
                "• Backend: PHP (Mô hình MVC thuần)",
                "• Database: MySQL (PDO kết nối an toàn)",
                "• Frontend: HTML5, CSS3, Bootstrap, Leaflet JS",
                "• AI Service: Trợ lý tư vấn lộ trình thông minh"
            ]
        },
        # Slide 5
        {
            "bg": DARK_BLUE,
            "title": "KIẾN TRÚC HỆ THỐNG (MVC)",
            "title_color": GOLD,
            "title_size": 36,
            "subtitle": "Quy trình xử lý khép kín, phân tách rõ ràng",
            "subtitle_color": WHITE,
            "subtitle_size": 16,
            "points": [
                "• Request: Khách hàng yêu cầu qua Router (index.php)",
                "• Controller: Nhận yêu cầu, điều hướng và xử lý logic",
                "• Model: Thực hiện truy vấn MySQL qua PDO",
                "• View: Render HTML/JS hiển thị bản đồ Leaflet"
            ]
        },
        # Slide 6
        {
            "bg": DARK_BLUE,
            "title": "CHỨC NĂNG 1: BẢN ĐỒ TƯƠNG TÁC",
            "title_color": GOLD,
            "title_size": 36,
            "subtitle": "Trực quan hóa hành trình bằng GPS",
            "subtitle_color": WHITE,
            "subtitle_size": 16,
            "points": [
                "• Thêm địa điểm: Click trực tiếp trên bản đồ số Leaflet",
                "• Tọa độ thực: Tự động bắt vĩ độ (Lat) và kinh độ (Lng)",
                "• Marker động: Phân loại điểm đi theo icon cá nhân"
            ]
        },
        # Slide 7
        {
            "bg": DARK_BLUE,
            "title": "CHỨC NĂNG 2: ALBUM KỶ NIỆM & CẢM XÚC",
            "title_color": GOLD,
            "title_size": 36,
            "subtitle": "Lưu trữ khoảnh khắc đẹp",
            "subtitle_color": WHITE,
            "subtitle_size": 16,
            "points": [
                "• Đa ảnh: Đăng tải nhiều ảnh kỷ niệm tại mỗi địa điểm",
                "• Cảm xúc: Đánh giá tâm trạng (Vui, Hào hứng, Bình yên...)",
                "• Tải trang nhanh: Đăng ảnh không reload nhờ Fetch API"
            ]
        },
        # Slide 8
        {
            "bg": DARK_BLUE,
            "title": "CHỨC NĂNG 3: CHUYẾN ĐI CHUNG (TRIPS)",
            "title_color": GOLD,
            "title_size": 36,
            "subtitle": "Đồng hành cùng bạn bè",
            "subtitle_color": WHITE,
            "subtitle_size": 16,
            "points": [
                "• Lập kế hoạch: Tạo hành trình chung (Start Date -> End Date)",
                "• Thành viên: Mời bạn bè cùng tham gia đóng góp địa điểm",
                "• Bản đồ nhóm: Hiển thị toàn bộ dấu chân của các thành viên"
            ]
        },
        # Slide 9
        {
            "bg": DARK_BLUE,
            "title": "CHỨC NĂNG 4: MẠNG XÃ HỘI THU NHỎ",
            "title_color": GOLD,
            "title_size": 36,
            "subtitle": "Gắn kết những người yêu xê dịch",
            "subtitle_color": WHITE,
            "subtitle_size": 16,
            "points": [
                "• Kết bạn: Tìm kiếm, gửi yêu cầu và duyệt bạn bè",
                "• Tương tác: Thả tim (Like), Bình luận (Comment) địa điểm",
                "• Tin nhắn riêng: Trò chuyện bảo mật thời gian thực"
            ]
        },
        # Slide 10
        {
            "bg": DARK_BLUE,
            "title": "CHỨC NĂNG 5: TRỢ LÝ AI THÔNG MINH",
            "title_color": GOLD,
            "title_size": 36,
            "subtitle": "Cá nhân hóa trải nghiệm du lịch",
            "subtitle_color": WHITE,
            "subtitle_size": 16,
            "points": [
                "• Gợi ý: Tư vấn điểm đến phù hợp sở thích người dùng",
                "• Lên lịch trình: Tạo kế hoạch chi tiết từng ngày tự động",
                "• Chatbot: Giải đáp thông tin thời tiết, ẩm thực địa phương"
            ]
        },
        # Slide 11
        {
            "bg": DARK_BLUE,
            "title": "PHÂN QUYỀN HỆ THỐNG",
            "title_color": GOLD,
            "title_size": 36,
            "subtitle": "Phân chia vai trò rõ ràng",
            "subtitle_color": WHITE,
            "subtitle_size": 16,
            "points": [
                "• Administrator: Quản trị hệ thống, cài đặt, khóa tài khoản",
                "• Moderator: Duyệt bài viết, quản lý báo cáo, hỗ trợ",
                "• User: Quản lý bản đồ cá nhân, chuyến đi và bạn bè",
                "• Privacy: Chế độ Công khai / Bạn bè / Riêng tư cho mỗi điểm"
            ]
        },
        # Slide 12
        {
            "bg": DARK_BLUE,
            "title": "BẢO MẬT HỆ THỐNG",
            "title_color": GOLD,
            "title_size": 36,
            "subtitle": "Bảo vệ thông tin người dùng",
            "subtitle_color": WHITE,
            "subtitle_size": 16,
            "points": [
                "• Mật khẩu: Mã hóa Bcrypt (password_hash) ngẫu nhiên muối",
                "• SQL Injection: Chống 100% bằng PDO Prepared Statements",
                "• XSS: Làm sạch dữ liệu đầu vào qua htmlspecialchars",
                "• Logs: Ghi lại hoạt động đăng nhập, IP người dùng"
            ]
        },
        # Slide 13
        {
            "bg": DARK_BLUE,
            "title": "DEMO & TRIỂN KHAI",
            "title_color": GOLD,
            "title_size": 36,
            "subtitle": "Sẵn sàng hoạt động thực tế",
            "subtitle_color": WHITE,
            "subtitle_size": 16,
            "points": [
                "• Địa chỉ: https://travelmap.page.gd",
                "• Hosting: Đã triển khai trực tiếp lên Internet",
                "• Tương thích: Tự động tối ưu giao diện trên Mobile & PC"
            ]
        },
        # Slide 14
        {
            "bg": DARK_BLUE,
            "title": "HƯỚNG PHÁT TRIỂN",
            "title_color": GOLD,
            "title_size": 36,
            "subtitle": "Kế hoạch nâng cấp tương lai",
            "subtitle_color": WHITE,
            "subtitle_size": 16,
            "points": [
                "• Ứng dụng di động: Đóng gói thành app Native hoàn chỉnh",
                "• Cloud Storage: Đẩy ảnh lên AWS S3 / Cloudinary tối ưu RAM",
                "• GPS Offline: Hỗ trợ lưu trữ hành trình khi mất kết nối mạng"
            ]
        },
        # Slide 15
        {
            "bg": DARK_BLUE,
            "title": "CẢM ƠN THẦY CÔ & HỘI ĐỒNG",
            "title_color": GOLD,
            "title_size": 40,
            "subtitle": "Q&A - Em xin phép nhận câu hỏi và đóng góp ý kiến",
            "subtitle_color": WHITE,
            "subtitle_size": 20,
            "points": []
        }
    ]
    
    for slide_data in slides_data:
        slide = prs.slides.add_slide(blank_slide_layout)
        apply_background(slide, slide_data["bg"])
        
        # Add Title Box
        title_box = slide.shapes.add_textbox(Inches(1.0), Inches(0.8), Inches(11.33), Inches(1.5))
        tf = title_box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        p = tf.paragraphs[0]
        p.text = slide_data["title"]
        p.alignment = PP_ALIGN.LEFT
        p.font.size = Pt(slide_data["title_size"])
        p.font.bold = True
        p.font.name = "Arial"
        p.font.color.rgb = slide_data["title_color"]
        
        # Add Subtitle Box if exists
        curr_top = 2.4
        if slide_data["subtitle"]:
            sub_box = slide.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.33), Inches(0.8))
            s_tf = sub_box.text_frame
            s_tf.word_wrap = True
            s_tf.margin_left = s_tf.margin_top = s_tf.margin_right = s_tf.margin_bottom = 0
            sp = s_tf.paragraphs[0]
            sp.text = slide_data["subtitle"]
            sp.font.size = Pt(slide_data["subtitle_size"])
            sp.font.italic = True
            sp.font.name = "Arial"
            sp.font.color.rgb = slide_data["subtitle_color"]
            curr_top = 3.0
            
        # Add Points Box
        if slide_data["points"]:
            points_box = slide.shapes.add_textbox(Inches(1.0), Inches(curr_top), Inches(11.33), Inches(4.0))
            p_tf = points_box.text_frame
            p_tf.word_wrap = True
            p_tf.margin_left = p_tf.margin_top = p_tf.margin_right = p_tf.margin_bottom = 0
            
            for i, pt in enumerate(slide_data["points"]):
                if i == 0:
                    pp = p_tf.paragraphs[0]
                else:
                    pp = p_tf.add_paragraph()
                pp.text = pt
                pp.font.size = Pt(22)
                pp.font.name = "Arial"
                pp.font.color.rgb = WHITE
                pp.space_after = Pt(14)
                
    # Save file to Desktop
    desktop_path = os.path.join(os.environ['USERPROFILE'], 'Desktop')
    
    desktop_options = [
        os.path.join(os.environ['USERPROFILE'], 'Desktop'),
        os.path.join(os.environ['USERPROFILE'], 'OneDrive', 'Desktop'),
        os.path.join(os.environ['USERPROFILE'], 'OneDrive - HUST', 'Desktop'),
        os.path.join(os.environ['USERPROFILE'], 'OneDrive - Hanoi University of Science and Technology', 'Desktop')
    ]
    
    target_dir = None
    for path in desktop_options:
        if os.path.exists(path):
            target_dir = path
            break
    if not target_dir:
        target_dir = os.environ['USERPROFILE']
        
    out_file = os.path.join(target_dir, 'Travelmap.pptx')
    prs.save(out_file)
    print(f"Presentation saved successfully to {out_file}")

if __name__ == "__main__":
    create_pptx()
