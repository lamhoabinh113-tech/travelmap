import os
import sys
import subprocess

# Check and install python-pptx if not present
try:
    import pptx
except ImportError:
    print("Installing python-pptx...")
    subprocess.run([sys.executable, "-m", "pip", "install", "python-pptx"], check=True)
    import pptx

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def apply_background(slide, color):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_header(slide, title_text, subtitle_text, gold_color, white_color):
    """Creates a standard modern header on slides to look like Canva templates"""
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(1.2))
    tf = title_box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
    
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.name = "Segoe UI"
    p.font.color.rgb = gold_color
    
    if subtitle_text:
        p2 = tf.add_paragraph()
        p2.text = subtitle_text
        p2.font.size = Pt(14)
        p2.font.italic = True
        p2.font.name = "Segoe UI"
        p2.font.color.rgb = white_color
        p2.space_before = Pt(4)

def add_gradient_card(slide, left, top, width, height, bg_color):
    """Adds a beautiful rounded rectangle shape to serve as a card background"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, 
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.color.rgb = RGBColor(46, 64, 87)
    shape.line.width = Pt(1)
    return shape

def create_pptx_premium():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]
    
    # Canva Color Scheme
    DARK_BLUE = RGBColor(15, 23, 42)      # Deep Navy Dark Background
    CARD_BLUE = RGBColor(30, 41, 59)      # Slate Blue Card Fill
    GOLD = RGBColor(245, 158, 11)         # Bright Amber / Gold Accent
    WHITE = RGBColor(248, 250, 252)       # Soft White body text
    MUTED_GRAY = RGBColor(148, 163, 184)  # Cool Grey body text
    
    conv_id = "b62215a2-c911-4a74-b42b-aa5a7f329653"
    brain_dir = os.path.join(os.environ['USERPROFILE'], '.gemini', 'antigravity', 'brain', conv_id)
    
    # Helper to check image existence
    def get_image_path(img_name):
        full_path = os.path.join(brain_dir, img_name)
        if os.path.exists(full_path):
            return full_path
        # Fallback to current folder just in case
        if os.path.exists(img_name):
            return img_name
        return None

    # SLIDE 1: Title Slide (Split design)
    slide1 = prs.slides.add_slide(blank_layout)
    apply_background(slide1, DARK_BLUE)
    
    # Left container for text
    tb1 = slide1.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(6.0), Inches(4.5))
    tf1 = tb1.text_frame
    tf1.word_wrap = True
    
    p = tf1.paragraphs[0]
    p.text = "TRAVEL MEMORY MAP"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.name = "Segoe UI"
    p.font.color.rgb = GOLD
    
    p2 = tf1.add_paragraph()
    p2.text = "Bản Đồ Kỷ Niệm Du Lịch Của Tôi"
    p2.font.size = Pt(28)
    p2.font.name = "Segoe UI"
    p2.font.color.rgb = WHITE
    p2.space_before = Pt(10)
    
    p3 = tf1.add_paragraph()
    p3.text = "\n\n• Sinh viên thực hiện: Lâm Hòa Bình\n• Công nghệ: PHP MVC & Leaflet JS"
    p3.font.size = Pt(16)
    p3.font.name = "Segoe UI"
    p3.font.color.rgb = MUTED_GRAY
    p3.space_before = Pt(20)

    # Right side: App Map Screenshot if exists
    map_img = get_image_path("media__1781079354989.png") or get_image_path("media__1781079276299.png")
    if map_img:
        slide1.shapes.add_picture(map_img, Inches(7.2), Inches(1.5), Inches(5.2), Inches(4.5))
    else:
        # Draw a beautiful graphic block as placeholder
        add_gradient_card(slide1, 7.2, 1.5, 5.2, 4.5, CARD_BLUE)
        tb_g = slide1.shapes.add_textbox(Inches(7.5), Inches(3.0), Inches(4.6), Inches(2.0))
        tb_g.text_frame.word_wrap = True
        gp = tb_g.text_frame.paragraphs[0]
        gp.text = "Bản đồ tương tác số\nLeaflet JS + PHP MVC"
        gp.alignment = PP_ALIGN.CENTER
        gp.font.size = Pt(24)
        gp.font.bold = True
        gp.font.color.rgb = GOLD

    # SLIDE 2: Problem & Solution (Card layout)
    slide2 = prs.slides.add_slide(blank_layout)
    apply_background(slide2, DARK_BLUE)
    add_header(slide2, "ĐẶT VẤN ĐỀ & GIẢI PHÁP", "Tại sao Travel Memory Map ra đời?", GOLD, WHITE)
    
    # 2 Cards
    add_gradient_card(slide2, 1.0, 2.0, 5.2, 4.3, CARD_BLUE)
    tb2_left = slide2.shapes.add_textbox(Inches(1.3), Inches(2.3), Inches(4.6), Inches(3.7))
    tf_left = tb2_left.text_frame
    tf_left.word_wrap = True
    lp1 = tf_left.paragraphs[0]
    lp1.text = "VẤN ĐỀ THỰC TẾ"
    lp1.font.bold = True
    lp1.font.size = Pt(22)
    lp1.font.color.rgb = GOLD
    
    lp2 = tf_left.add_paragraph()
    lp2.text = "\n• Lưu trữ ảnh du lịch rời rạc trên điện thoại\n• Thiếu gắn kết không gian và lộ trình địa lý thực tế\n• Khó tìm lại khoảnh khắc theo vị trí"
    lp2.font.size = Pt(16)
    lp2.font.color.rgb = WHITE
    lp2.space_before = Pt(10)
    
    add_gradient_card(slide2, 7.0, 2.0, 5.2, 4.3, CARD_BLUE)
    tb2_right = slide2.shapes.add_textbox(Inches(7.3), Inches(2.3), Inches(4.6), Inches(3.7))
    tf_right = tb2_right.text_frame
    tf_right.word_wrap = True
    rp1 = tf_right.paragraphs[0]
    rp1.text = "GIẢI PHÁP"
    rp1.font.bold = True
    rp1.font.size = Pt(22)
    rp1.font.color.rgb = GOLD
    
    rp2 = tf_right.add_paragraph()
    rp2.text = "\n• Số hóa hành trình bằng bản đồ Leaflet JS\n• Ghim ảnh kỷ niệm và ghi chú cảm xúc vào tọa độ GPS thực\n• Trực quan hóa toàn bộ dấu chân du lịch"
    rp2.font.size = Pt(16)
    rp2.font.color.rgb = WHITE
    rp2.space_before = Pt(10)

    # SLIDE 3: Objectives (3 Cards)
    slide3 = prs.slides.add_slide(blank_layout)
    apply_background(slide3, DARK_BLUE)
    add_header(slide3, "MỤC TIÊU ĐỀ TÀI", "Định hướng phát triển hệ thống", GOLD, WHITE)
    
    card_w = 3.6
    for i, (title, desc) in enumerate([
        ("SỐ HÓA HÀNH TRÌNH", "Lưu trữ tọa độ GPS, hình ảnh kỷ niệm, ghi chú và dòng thời gian các chuyến đi."),
        ("KẾT NỐI BẠN BÈ", "Chia sẻ chuyến đi chung nhóm, tương tác thả tim, bình luận và chat trực tuyến."),
        ("TƯƠNG TÁC AI", "Tích hợp chatbot AI gợi ý điểm đến, lên lịch trình du lịch thông minh cá nhân hóa.")
    ]):
        left_pos = 0.8 + i * 4.1
        add_gradient_card(slide3, left_pos, 2.2, card_w, 4.0, CARD_BLUE)
        tb_c = slide3.shapes.add_textbox(Inches(left_pos + 0.2), Inches(2.5), Inches(card_w - 0.4), Inches(3.4))
        tf_c = tb_c.text_frame
        tf_c.word_wrap = True
        cp1 = tf_c.paragraphs[0]
        cp1.text = f"0{i+1}\n\n{title}"
        cp1.font.bold = True
        cp1.font.size = Pt(20)
        cp1.font.color.rgb = GOLD
        
        cp2 = tf_c.add_paragraph()
        cp2.text = f"\n{desc}"
        cp2.font.size = Pt(15)
        cp2.font.color.rgb = WHITE

    # SLIDE 4: Technologies Used (Left points, Right Diagram)
    slide4 = prs.slides.add_slide(blank_layout)
    apply_background(slide4, DARK_BLUE)
    add_header(slide4, "CÔNG NGHỆ SỬ DỤNG", "Nền tảng tối ưu hiệu năng", GOLD, WHITE)
    
    tb4 = slide4.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(5.8), Inches(4.5))
    tf4 = tb4.text_frame
    tf4.word_wrap = True
    for tech, desc in [
        ("Backend PHP MVC", "Mô hình sạch, xử lý nhanh chóng ở server."),
        ("MySQL & PDO", "Kết nối an toàn, bảo mật cao chống SQL Injection."),
        ("Leaflet JS", "Bản đồ số mượt mà, định vị GPS chính xác."),
        ("Fetch API (Ajax)", "Tương tác đăng ảnh mượt, không reload trang.")
    ]:
        p = tf4.add_paragraph()
        p.text = f"• {tech}: {desc}"
        p.font.size = Pt(16)
        p.font.color.rgb = WHITE
        p.space_after = Pt(14)
        
    client_server_img = get_image_path("client_server_diagram_1781060963217.png")
    if client_server_img:
        slide4.shapes.add_picture(client_server_img, Inches(7.0), Inches(2.0), Inches(5.5), Inches(4.2))
    else:
        add_gradient_card(slide4, 7.0, 2.0, 5.5, 4.2, CARD_BLUE)

    # SLIDE 5: MVC Architecture (Flow chart slide)
    slide5 = prs.slides.add_slide(blank_layout)
    apply_background(slide5, DARK_BLUE)
    add_header(slide5, "KIẾN TRÚC HỆ THỐNG (MVC)", "Phân tách vai trò xử lý dữ liệu và giao diện", GOLD, WHITE)
    
    mvc_img = get_image_path("mvc_travel_diagram_1781062196390.png") or get_image_path("mvc_diagram_1781060943701.png")
    if mvc_img:
        slide5.shapes.add_picture(mvc_img, Inches(0.8), Inches(1.8), Inches(7.0), Inches(4.8))
        tb5 = slide5.shapes.add_textbox(Inches(8.2), Inches(2.0), Inches(4.5), Inches(4.2))
        tf5 = tb5.text_frame
        tf5.word_wrap = True
        for name, role in [
            ("Router", "Định tuyến URL thân thiện (index.php)"),
            ("Controller", "Điều phối logic nghiệp vụ chính"),
            ("Model", "Truy vấn MySQL an toàn qua PDO"),
            ("View", "Hiển thị HTML/CSS/JS bản đồ Leaflet")
        ]:
            p = tf5.add_paragraph()
            p.text = f"✔ {name}\n  {role}"
            p.font.size = Pt(14)
            p.font.color.rgb = WHITE
            p.space_after = Pt(10)
    else:
        # Generic split layout if no diagram
        tb5 = slide5.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(11.33), Inches(4.5))
        tf5 = tb5.text_frame
        tf5.word_wrap = True
        p = tf5.paragraphs[0]
        p.text = "Quy trình MVC trong ứng dụng:"
        p.font.bold = True
        p.font.size = Pt(22)
        p.font.color.rgb = GOLD
        # Add MVC flow text
        
    # SLIDE 6: Interactive Map
    slide6 = prs.slides.add_slide(blank_layout)
    apply_background(slide6, DARK_BLUE)
    add_header(slide6, "CHỨC NĂNG 1: BẢN ĐỒ TƯƠNG TÁC", "Trực quan hóa điểm đến GPS", GOLD, WHITE)
    
    tb6 = slide6.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(5.8), Inches(4.5))
    tf6 = tb6.text_frame
    tf6.word_wrap = True
    for point in [
        "• Ghim vị trí trực tiếp bằng cách Click bản đồ",
        "• Tự động thu thập Kinh độ & Vĩ độ thực (GPS)",
        "• Phân biệt loại địa điểm thông qua các Icon Marker",
        "• Tìm kiếm địa điểm nhanh chóng thông qua địa danh"
    ]:
        p = tf6.add_paragraph()
        p.text = point
        p.font.size = Pt(18)
        p.font.color.rgb = WHITE
        p.space_after = Pt(16)
        
    map_scr = get_image_path("media__1781079354989.png")
    if map_scr:
        slide6.shapes.add_picture(map_scr, Inches(7.0), Inches(2.0), Inches(5.5), Inches(4.2))
    else:
        add_gradient_card(slide6, 7.0, 2.0, 5.5, 4.2, CARD_BLUE)

    # SLIDE 7: Album and Emotion
    slide7 = prs.slides.add_slide(blank_layout)
    apply_background(slide7, DARK_BLUE)
    add_header(slide7, "CHỨC NĂNG 2: ALBUM KỶ NIỆM & CẢM XÚC", "Lưu lại những khoảnh khắc và tâm trạng", GOLD, WHITE)
    
    tb7 = slide7.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(5.8), Inches(4.5))
    tf7 = tb7.text_frame
    tf7.word_wrap = True
    for point in [
        "• Tải nhiều ảnh cùng lúc lên từng địa điểm",
        "• Gắn nhãn cảm xúc: Hào hứng, Bình yên, Vui vẻ...",
        "• Tương tác mượt: Đăng ảnh không reload trang",
        "• Album dạng lưới hiển thị ảnh tràn viền tối ưu"
    ]:
        p = tf7.add_paragraph()
        p.text = point
        p.font.size = Pt(18)
        p.font.color.rgb = WHITE
        p.space_after = Pt(16)
        
    album_img = get_image_path("media__1781079289734.png")
    if album_img:
        slide7.shapes.add_picture(album_img, Inches(7.0), Inches(2.0), Inches(5.5), Inches(4.2))
    else:
        add_gradient_card(slide7, 7.0, 2.0, 5.5, 4.2, CARD_BLUE)

    # SLIDE 8: Shared Trips
    slide8 = prs.slides.add_slide(blank_layout)
    apply_background(slide8, DARK_BLUE)
    add_header(slide8, "CHỨC NĂNG 3: HÀNH TRÌNH CHUNG (TRIPS)", "Lên kế hoạch và đồng hành cùng nhóm bạn", GOLD, WHITE)
    
    tb8 = slide8.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(5.8), Inches(4.5))
    tf8 = tb8.text_frame
    tf8.word_wrap = True
    for point in [
        "• Tạo hành trình chung: Tên, Ghi chú, Ngày bắt đầu & kết thúc",
        "• Mời bạn bè cùng tham gia đóng góp điểm đi trên bản đồ",
        "• Hiển thị lộ trình nhóm bằng các đường vẽ (polyline) nối tiếp",
        "• Quản lý danh sách thành viên chuyến đi dễ dàng"
    ]:
        p = tf8.add_paragraph()
        p.text = point
        p.font.size = Pt(18)
        p.font.color.rgb = WHITE
        p.space_after = Pt(16)
        
    trip_img = get_image_path("media__1781077244500.png")
    if trip_img:
        slide8.shapes.add_picture(trip_img, Inches(7.0), Inches(2.0), Inches(5.5), Inches(4.2))
    else:
        add_gradient_card(slide8, 7.0, 2.0, 5.5, 4.2, CARD_BLUE)

    # SLIDE 9: Social network interaction
    slide9 = prs.slides.add_slide(blank_layout)
    apply_background(slide9, DARK_BLUE)
    add_header(slide9, "CHỨC NĂNG 4: MẠNG XÃ HỘI DU LỊCH", "Gắn kết cộng đồng yêu khám phá", GOLD, WHITE)
    
    tb9 = slide9.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(5.8), Inches(4.5))
    tf9 = tb9.text_frame
    tf9.word_wrap = True
    for point in [
        "• Gửi và phê duyệt lời mời kết bạn",
        "• Thả tim (Like) & Bình luận (Comment) tại từng địa điểm",
        "• Nhắn tin riêng tư trò chuyện trực tuyến",
        "• News Feed cập nhật kỷ niệm mới của bạn bè"
    ]:
        p = tf9.add_paragraph()
        p.text = point
        p.font.size = Pt(18)
        p.font.color.rgb = WHITE
        p.space_after = Pt(16)
        
    feed_img = get_image_path("media__1781071852596.png")
    if feed_img:
        slide9.shapes.add_picture(feed_img, Inches(7.0), Inches(2.0), Inches(5.5), Inches(4.2))
    else:
        add_gradient_card(slide9, 7.0, 2.0, 5.5, 4.2, CARD_BLUE)

    # SLIDE 10: AI Assistant
    slide10 = prs.slides.add_slide(blank_layout)
    apply_background(slide10, DARK_BLUE)
    add_header(slide10, "CHỨC NĂNG 5: TRỢ LÝ TRÍ TUỆ NHÂN TẠO", "Cá nhân hóa lịch trình thông minh", GOLD, WHITE)
    
    tb10 = slide10.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(5.8), Inches(4.5))
    tf10 = tb10.text_frame
    tf10.word_wrap = True
    for point in [
        "• Chatbot tư vấn trực tiếp điểm đến thời gian thực",
        "• Tạo lịch trình du lịch tự động theo số ngày yêu cầu",
        "• Phân tích sở thích từ lịch sử điểm đi để gợi ý vị trí",
        "• Tương tác nhanh, giao diện tối giản, sinh động"
    ]:
        p = tf10.add_paragraph()
        p.text = point
        p.font.size = Pt(18)
        p.font.color.rgb = WHITE
        p.space_after = Pt(16)
        
    ai_img = get_image_path("ai_chat_mockup_1781062241590.png")
    if ai_img:
        slide10.shapes.add_picture(ai_img, Inches(7.0), Inches(2.0), Inches(5.5), Inches(4.2))
    else:
        add_gradient_card(slide10, 7.0, 2.0, 5.5, 4.2, CARD_BLUE)

    # SLIDE 11: Permissions and Settings
    slide11 = prs.slides.add_slide(blank_layout)
    apply_background(slide11, DARK_BLUE)
    add_header(slide11, "PHÂN QUYỀN HỆ THỐNG & RIÊNG TƯ", "Đảm bảo an toàn thông tin và quyền hạn", GOLD, WHITE)
    
    # 2 Cards Side-by-Side
    add_gradient_card(slide11, 1.0, 2.0, 5.2, 4.3, CARD_BLUE)
    tb11_l = slide11.shapes.add_textbox(Inches(1.3), Inches(2.3), Inches(4.6), Inches(3.7))
    tf11_l = tb11_l.text_frame
    tf11_l.word_wrap = True
    p = tf11_l.paragraphs[0]
    p.text = "VAI TRÒ (ROLES)"
    p.font.bold = True
    p.font.size = Pt(22)
    p.font.color.rgb = GOLD
    p_desc = tf11_l.add_paragraph()
    p_desc.text = "\n• Admin: Cấu hình hệ thống, khóa user, theo dõi hoạt động.\n• Moderator: Duyệt báo cáo, kiểm duyệt hình ảnh vi phạm.\n• User: Sở hữu bản đồ cá nhân, thêm hành trình du lịch."
    p_desc.font.size = Pt(16)
    p_desc.font.color.rgb = WHITE
    p_desc.space_before = Pt(10)
    
    add_gradient_card(slide11, 7.0, 2.0, 5.2, 4.3, CARD_BLUE)
    tb11_r = slide11.shapes.add_textbox(Inches(7.3), Inches(2.3), Inches(4.6), Inches(3.7))
    tf11_r = tb11_r.text_frame
    tf11_r.word_wrap = True
    p_r = tf11_r.paragraphs[0]
    p_r.text = "QUYỀN RIÊNG TƯ (PRIVACY)"
    p_r.font.bold = True
    p_r.font.size = Pt(22)
    p_r.font.color.rgb = GOLD
    p_r_desc = tf11_r.add_paragraph()
    p_r_desc.text = "\n• Công khai (Public): Mọi người có thể xem trên News Feed.\n• Bạn bè (Friends): Chỉ bạn bè đã duyệt mới được xem.\n• Riêng tư (Private): Chỉ chủ sở hữu nhìn thấy trên bản đồ."
    p_r_desc.font.size = Pt(16)
    p_r_desc.font.color.rgb = WHITE
    p_r_desc.space_before = Pt(10)

    # SLIDE 12: Security Mechanisms
    slide12 = prs.slides.add_slide(blank_layout)
    apply_background(slide12, DARK_BLUE)
    add_header(slide12, "BẢO MẬT DỮ LIỆU", "Giải pháp bảo vệ tài khoản và thông tin", GOLD, WHITE)
    
    tb12 = slide12.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(5.8), Inches(4.5))
    tf12 = tb12.text_frame
    tf12.word_wrap = True
    for sec in [
        "• Mật khẩu: Mã hóa BCRYPT an toàn vượt trội MD5",
        "• SQL Injection: Chặn đứng bằng PDO Prepared Statements",
        "• XSS Clean: Lọc sạch HTML/Mã độc đầu vào (htmlspecialchars)",
        "• Login Logs: Giám sát lịch sử đăng nhập, thiết bị và IP"
    ]:
        p = tf12.add_paragraph()
        p.text = sec
        p.font.size = Pt(18)
        p.font.color.rgb = WHITE
        p.space_after = Pt(18)
        
    db_img = get_image_path("db_schema_travel_1781062265135.png")
    if db_img:
        slide12.shapes.add_picture(db_img, Inches(7.0), Inches(2.0), Inches(5.5), Inches(4.2))
    else:
        add_gradient_card(slide12, 7.0, 2.0, 5.5, 4.2, CARD_BLUE)

    # SLIDE 13: Demo and Hosting (Responsive preview if possible)
    slide13 = prs.slides.add_slide(blank_layout)
    apply_background(slide13, DARK_BLUE)
    add_header(slide13, "LIVE DEMO & TRIỂN KHAI", "Hệ thống sẵn sàng vận hành thực tế", GOLD, WHITE)
    
    # Left Card
    add_gradient_card(slide13, 1.0, 2.0, 5.2, 4.3, CARD_BLUE)
    tb13_l = slide13.shapes.add_textbox(Inches(1.3), Inches(2.3), Inches(4.6), Inches(3.7))
    tf13_l = tb13_l.text_frame
    tf13_l.word_wrap = True
    lp = tf13_l.paragraphs[0]
    lp.text = "THÔNG TIN TRIỂN KHAI"
    lp.font.bold = True
    lp.font.size = Pt(22)
    lp.font.color.rgb = GOLD
    lp_desc = tf13_l.add_paragraph()
    lp_desc.text = "\n• Tên miền: travelmap.page.gd\n• Triển khai: Hosting Internet trực tiếp\n• Thiết lập: SSL an toàn (HTTPS)\n• Cơ sở dữ liệu: Đồng bộ đám mây"
    lp_desc.font.size = Pt(16)
    lp_desc.font.color.rgb = WHITE
    lp_desc.space_before = Pt(10)
    
    # Right Card
    add_gradient_card(slide13, 7.0, 2.0, 5.2, 4.3, CARD_BLUE)
    tb13_r = slide13.shapes.add_textbox(Inches(7.3), Inches(2.3), Inches(4.6), Inches(3.7))
    tf13_r = tb13_r.text_frame
    tf13_r.word_wrap = True
    rp = tf13_r.paragraphs[0]
    rp.text = "RESPONSIVE MOBILE"
    rp.font.bold = True
    rp.font.size = Pt(22)
    rp.font.color.rgb = GOLD
    rp_desc = tf13_r.add_paragraph()
    rp_desc.text = "\n• Thiết kế Web tương thích mọi thiết bị di động\n• Thao tác bản đồ chạm kéo (Touch & Drag) mượt mà\n• Hỗ trợ cài đặt nhanh lên điện thoại qua PWA"
    rp_desc.font.size = Pt(16)
    rp_desc.font.color.rgb = WHITE
    rp_desc.space_before = Pt(10)

    # SLIDE 14: Future roadmap
    slide14 = prs.slides.add_slide(blank_layout)
    apply_background(slide14, DARK_BLUE)
    add_header(slide14, "ĐỊNH HƯỚNG PHÁT TRIỂN", "Giải pháp nâng cấp và mở rộng", GOLD, WHITE)
    
    card_w = 3.6
    for i, (num, title, desc) in enumerate([
        ("01", "APP DI ĐỘNG", "Đóng gói ứng dụng thành ứng dụng di động gốc (Flutter/React Native) để tải lên AppStore/GooglePlay."),
        ("02", "CLOUD STORAGE", "Tích hợp lưu trữ hình ảnh thông qua AWS S3 hoặc Cloudinary để gánh băng thông và lưu lượng."),
        ("03", "MARKER CLUSTERING", "Tự động gom nhóm các điểm đi ở gần nhau để bản đồ luôn hiển thị gọn gàng, trực quan.")
    ]):
        left_pos = 0.8 + i * 4.1
        add_gradient_card(slide14, left_pos, 2.2, card_w, 4.0, CARD_BLUE)
        tb_c = slide14.shapes.add_textbox(Inches(left_pos + 0.2), Inches(2.5), Inches(card_w - 0.4), Inches(3.4))
        tf_c = tb_c.text_frame
        tf_c.word_wrap = True
        cp1 = tf_c.paragraphs[0]
        cp1.text = f"{num}\n\n{title}"
        cp1.font.bold = True
        cp1.font.size = Pt(20)
        cp1.font.color.rgb = GOLD
        
        cp2 = tf_c.add_paragraph()
        cp2.text = f"\n{desc}"
        cp2.font.size = Pt(15)
        cp2.font.color.rgb = WHITE

    # SLIDE 15: Q&A - Quyết Thuyết Trình Thành Công
    slide15 = prs.slides.add_slide(blank_layout)
    apply_background(slide15, DARK_BLUE)
    
    tb15 = slide15.shapes.add_textbox(Inches(1.5), Inches(2.2), Inches(10.3), Inches(4.0))
    tf15 = tb15.text_frame
    tf15.word_wrap = True
    
    p = tf15.paragraphs[0]
    p.text = "XIN CẢM ƠN THẦY CÔ & HỘI ĐỒNG!"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.name = "Segoe UI"
    p.font.color.rgb = GOLD
    
    p2 = tf15.add_paragraph()
    p2.text = "\nEm xin lắng nghe các câu hỏi và đóng góp ý kiến để hoàn thiện đề tài."
    p2.alignment = PP_ALIGN.CENTER
    p2.font.size = Pt(20)
    p2.font.italic = True
    p2.font.name = "Segoe UI"
    p2.font.color.rgb = WHITE
    p2.space_before = Pt(20)

    # Save to user desktop
    desktop_options = [
        os.path.join(os.environ['USERPROFILE'], 'Desktop'),
        os.path.join(os.environ['USERPROFILE'], 'OneDrive', 'Desktop'),
        os.path.join(os.environ['USERPROFILE'], 'OneDrive - HUST', 'Desktop'),
        os.path.join(os.environ['USERPROFILE'], 'OneDrive - Hanoi University of Science and Technology', 'Desktop')
    ]
    
    target_dir = None
    for path in desktop_options:
        if os.path.exists(path):
            target_dir = path
            break
    if not target_dir:
        target_dir = os.environ['USERPROFILE']
        
    out_file = os.path.join(target_dir, 'Travelmap_Premium.pptx')
    prs.save(out_file)
    print(f"Premium Canva-style Presentation saved successfully to {out_file}")

if __name__ == "__main__":
    create_pptx_premium()
